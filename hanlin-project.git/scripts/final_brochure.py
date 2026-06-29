#!/usr/bin/env python3
"""
创建最终版宣传册（混合配图策略）
- 实景图：ComfyUI 生成校园场景图
- 数据可视化：ComfyUI 生成图表与图标
"""

from docx import Document
from pptx import Presentation
from pathlib import Path


def create_final_brochure():
    """创建最终版宣传册（DOCX + PPT）"""
    
    print("=" * 60)
    print("创建最终版宣传册")
    print("=" * 60)
    
    # Step 1: ComfyUI 生成的图像列表
    generated_images = {
        "school_exterior": "mcp_generated_00009_.png",  # PPCHA 校园场景图
        "infographic_charts": "mcp_generated_00010_.png",  # 数据可视化图表
        "cover_background": "mcp_generated_00006_.png"   # 封面背景图
    }
    
    print("\n🎨 ComfyUI 已生成图像:")
    for name, img_path in generated_images.items():
        print(f"   - {name}: {img_path}")
    
    # Step 2: Create DOCX with images
    doc = Document()
    
    # Slide 1: Cover with school exterior image
    title = "中华书院国际教育合作项目说明书"
    subtitle = "华侨生 985/211 升学项目（校长版）"
    
    doc.add_heading(title, level=1)
    
    p = doc.add_paragraph()
    run = p.add_run(subtitle)
    run.font.size = 14
    
    # Add ComfyUI-generated school exterior image to cover
    try:
        from docx.shared import Inches
        
        school_img_path = f"/root/.openclaw/workspace/EducationStudio/outputs/{generated_images['school_exterior']}"
        
        if Path(school_img_path).exists():
            logo = doc.add_picture(
                school_img_path, 
                width=Inches(6)
            )
            print(f"\n✅ DOCX 封面图已嵌入：{school_img_path}")
            
    except Exception as e:
        print(f"⚠️  DOCX 图像添加失败：{e}")
    
    # Slide 2: Why Partner? with infographic chart
    doc.add_heading("为什么选择中华书院？", level=1)
    
    p = doc.add_paragraph()
    run = p.add_run("""
📊 **三大核心优势：**

1. ✅ 升学保障 - 双向承诺书制度，未录取即退费
2. ✅ 多元路径 - 985/211 + 海外名校双通道  
3. ✅ 文化传承 - 中华文化与国际视野并重""")
    
    run.font.size = 10
    
    # Add ComfyUI-generated infographic chart
    try:
        chart_path = f"/root/.openclaw/workspace/EducationStudio/outputs/{generated_images['infographic_charts']}"
        
        if Path(chart_path).exists():
            chart = doc.add_picture(
                chart_path, 
                width=Inches(7)
            )
            print(f"✅ 数据图表已嵌入：{chart_path}")
            
    except Exception as e:
        print(f"⚠️  图表添加失败：{e}")
    
    # Save DOCX
    output_path = "/root/.openclaw/workspace/EducationStudio/outputs/中华书院_华侨生保录宣传册_校长版_final.docx"
    doc.save(output_path)
    
    print(f"\n✅ DOCX 已保存到：{output_path}")
    
    # Step 3: Create PPT with images
    prs = Presentation()
    
    # Slide 1: Cover
    slide_layout = prs.slide_layouts[0]  # Title only
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    content_box = slide.shapes.add_textbox(
        left=Inches(0.5),
        top=Inches(2.0),
        width=Inches(6.0),
        height=Inches(4.0)
    )
    
    title_shape.text = f"{title}\n{subtitle}"
    
    # Add ComfyUI-generated school exterior image to PPT cover
    try:
        from pptx.util import Inches
        
        if Path(school_img_path).exists():
            shape = slide.shapes.add_picture(
                school_img_path,
                left=Inches(0.5),
                top=Inches(2.0),
                width=Inches(6.0)
            )
            
    except Exception as e:
        print(f"⚠️  PPT 封面图添加失败：{e}")
    
    # Slide 2: Why Partner? with infographic
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    content_box = slide.shapes.add_textbox(
        left=Inches(0.5),
        top=Inches(2.0),
        width=Inches(6.0),
        height=Inches(4.0)
    )
    
    title_shape.text = "为什么选择中华书院？"
    
    tf = content_box.text_frame
    for point in [
        "- 📉 生源竞争持续加剧",
        "- 🎯 家长需求升级转变", 
        "- 💡 发展路径亟待拓宽",
        ""
    ]:
        if point:
            p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
            p.text = point
    
    # Add ComfyUI-generated infographic chart to PPT
    try:
        from pptx.util import Inches
        
        if Path(chart_path).exists():
            shape = slide.shapes.add_picture(
                chart_path,
                left=Inches(7.0),
                top=Inches(1.0),
                width=Inches(5.0)
            )
            print(f"✅ PPT 数据图表已嵌入：{chart_path}")
            
    except Exception as e:
        print(f"⚠️  PPT 图表添加失败：{e}")
    
    # Save PPTX
    ppt_output_path = "/root/.openclaw/workspace/EducationStudio/outputs/中华书院_华侨生保录宣传册_校长版_final.pptx"
    prs.save(ppt_output_path)
    
    print(f"\n✅ PPTX 已保存到：{ppt_output_path}")
    
    # Step 4: Summary
    print("\n" + "=" * 60)
    print("📊 最终版宣传册生成完成！")
    print("=" * 60)
    print(f"\n📁 输出文件:")
    print(f"   📄 DOCX: {output_path}")
    print(f"   📊 PPTX: {ppt_output_path}")


if __name__ == "__main__":
    create_final_brochure()
