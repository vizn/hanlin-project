#!/usr/bin/env python3
"""Convert Markdown brochure/project description to branded PPT."""

import sys, json, re, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def load_brand(path):
    if os.path.exists(path):
        with open(path) as f:
            return json.load(f)
    return {"colors":{"primary":"#2196F3","secondary":"#FF9800","accent":"#E3F2FD","text_dark":"#333333","white":"#FFFFFF"},"fonts":{"heading":"Microsoft YaHei","body":"Microsoft YaHei"},"ppt":{"slide_width_cm":33.867,"slide_height_cm":19.05,"title_font_size":36,"subtitle_font_size":20,"body_font_size":16,"accent_bar":True}}

def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def create_slide(prs, title_text, body_lines, brand, level=1):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = hex_to_rgb(brand["colors"]["white"])
    W = prs.slide_width
    H = prs.slide_height
    primary = hex_to_rgb(brand["colors"]["primary"])
    secondary = hex_to_rgb(brand["colors"]["secondary"])
    text_dark = hex_to_rgb(brand["colors"]["text_dark"])
    accent = hex_to_rgb(brand["colors"]["accent"])
    pcfg = brand.get("ppt", {})
    title_size = pcfg.get("subtitle_font_size" if level > 1 else "title_font_size", 36)
    body_size = pcfg.get("body_font_size", 16)

    if pcfg.get("accent_bar", True):
        bar = slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(80000), H)  # rectangle
        bar.fill.solid()
        bar.fill.fore_color.rgb = primary
        bar.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Cm(2), Cm(1.5), Cm(28), Cm(3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(title_size)
    p.font.bold = True
    p.font.color.rgb = primary

    # Divider line
    div = slide.shapes.add_shape(1, Cm(2), Cm(4.5), Cm(6), Emu(30000))
    div.fill.solid()
    div.fill.fore_color.rgb = secondary
    div.line.fill.background()

    # Body
    txBox2 = slide.shapes.add_textbox(Cm(2), Cm(5.5), Cm(28), Cm(12))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, line in enumerate(body_lines):
        line = line.strip()
        if not line:
            continue
        if line.startswith("- ") or line.startswith("* "):
            line = "  " + line
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = line
        p.font.size = Pt(body_size)
        p.font.color.rgb = text_dark
        p.space_after = Pt(6)

def md_to_ppt(md_path, output_path, brand_path=None):
    brand_path = brand_path or os.path.join(os.path.dirname(__file__), "brand_template.json")
    brand = load_brand(brand_path)
    prs = Presentation()
    prs.slide_width = Cm(brand.get("ppt",{}).get("slide_width_cm", 33.867))
    prs.slide_height = Cm(brand.get("ppt",{}).get("slide_height_cm", 19.05))

    with open(md_path, encoding="utf-8") as f:
        lines = f.readlines()

    slides = []
    current_title = None
    current_body = []
    current_level = 1

    for line in lines:
        line = line.rstrip()
        if line.startswith("# "):
            if current_title:
                slides.append((current_title, current_body, current_level))
            current_title = line[2:]
            current_body = []
            current_level = 1
        elif line.startswith("## "):
            if current_title:
                slides.append((current_title, current_body, current_level))
            current_title = line[3:]
            current_body = []
            current_level = 2
        elif line.startswith("---"):
            if current_title:
                slides.append((current_title, current_body, current_level))
            current_title = None
            current_body = []
        else:
            if line.strip():
                current_body.append(line)
    if current_title:
        slides.append((current_title, current_body, current_level))

    for title, body, level in slides:
        create_slide(prs, title, body, brand, level)

    prs.save(output_path)
    print(f"PPT saved: {output_path}")

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: md2ppt.py <input.md> <output.pptx> [brand_template.json]")
        sys.exit(1)
    md_to_ppt(sys.argv[1], sys.argv[2], sys.argv[3] if len(sys.argv) > 3 else None)
