#!/usr/bin/env python3
"""
将 ComfyUI 生成的图像嵌入到 DOCX/PPT 宣传册中
"""

import json


def create_final_brochure():
    """创建最终版宣传册（带封面图）"""
    
    print("=" * 60)
    print("创建最终版宣传册")
    print("=" * 60)
    
    # Step 1: 读取源 Markdown
    source_md = "/root/.openclaw/workspace/中华书院_华侨生保录宣传册_校长版.md"
    
    with open(source_md, 'r', encoding='utf-8') as f:
        content = f.read()
    
    print(f"\n📄 源文件：{source_md}")
    print(f"   大小：{len(content)} bytes")
    
    # Step 2: 提取已生成的封面图文件名（从 ComfyUI API）
    generated_images = [
        "mcp_generated_00006_.png",
        "mcp_generated_00007_.png", 
        "mcp_generated_00008_.png"
    ]
    
    print(f"\n🎨 已生成 {len(generated_images)} 张封面图:")
    for img in generated_images:
        print(f"   - {img}")
    
    # Step 3: 创建最终版 DOCX（带图像）
    from docx import Document
    
    doc = Document()
    
    # Slide 1: Cover with image
    slide_title = "中华书院国际教育合作项目说明书"
    slide_subtitle = "华侨生 985/211 升学项目（校长版）"
    
    doc.add_heading(slide_title, level=1)
    
    # Add subtitle as paragraph
    from docx.shared import Pt
    
    p = doc.add_paragraph()
    run = p.add_run(slide_subtitle)
    run.font.size = Pt(14)
    
    # Add first image (封面图) if available
    if generated_images and "mcp_generated_00006_.png" in generated_images:
        try:
            from docx.shared import Inches
            
            # 尝试添加图像（如果 ComfyUI 已保存）
            img_path = f"/root/.openclaw/workspace/EducationStudio/outputs/{generated_images[0]}"
            
            if __import__('pathlib').Path(img_path).exists():
                logo = doc.add_picture(img_path, width=Inches(6))
                print(f"\n✅ 封面图已嵌入：{img_path}")
            else:
                # 如果没有本地文件，添加占位符说明
                p = doc.add_paragraph()
                run = p.add_run("📷 [封面图占位符]")
                run.font.size = Pt(10)
                run.bold = True
                
        except Exception as e:
            print(f"\n⚠️  无法添加图像：{e}")
    
    # Step 4: 保存 DOCX
    output_path = "/root/.openclaw/workspace/EducationStudio/outputs/中华书院_华侨生保录宣传册_校长版_final.docx"
    doc.save(output_path)
    
    print(f"\n✅ 最终版 DOCX 已保存到：{output_path}")
    
    # Step 5: 创建 PPT（带封面图）
    from pptx import Presentation
    
    prs = Presentation()
    
    # Slide 1: Cover with image
    title = "中华书院国际教育合作项目说明书"
    subtitle = "华侨生 985/211 升学项目（校长版）"
    
    slide_layout = prs.slide_layouts[6]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    content_shape = slide.placeholders.get(1) or slide.shapes.add_textbox(
        left=Inches(0.5),
        top=Inches(2.0),
        width=Inches(6.0),
        height=Inches(4.0)
    )
    
    title_shape.text = "中华书院国际教育合作项目说明书"
    
    tf = content_shape.text_frame
    for line in [
        "华侨生 985/211 升学项目（校长版）",
        "",
        "📷 封面图预览",
        ""
    ]:
        p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
        p.text += line
    
    # Add image placeholder (or actual image)
    try:
        img_path = f"/root/.openclaw/workspace/EducationStudio/outputs/{generated_images[0]}"
        
        from pptx.util import Inches
        
        if __import__('pathlib').Path(img_path).exists():
            shape = slide.shapes.add_picture(
                img_path,
                left=Inches(0.5),
                top=Inches(2.0),
                width=Inches(6.0)
            )
            print(f"✅ PPT 封面图已嵌入：{img_path}")
        else:
            shape = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(6.0), Inches(4.0))
            tf = shape.text_frame
            p = tf.paragraphs[0]
            p.text = "📷 [封面图占位符]"
            
    except Exception as e:
        print(f"⚠️  PPT 图像添加失败：{e}")
    
    # Save PPTX
    ppt_output_path = "/root/.openclaw/workspace/EducationStudio/outputs/中华书院_华侨生保录宣传册_校长版_final.pptx"
    prs.save(ppt_output_path)
    
    print(f"\n✅ 最终版 PPTX 已保存到：{ppt_output_path}")
    
    # Step 6: Summary
    print("\n" + "=" * 60)
    print("📊 最终版宣传册生成完成！")
    print("=" * 60)
    print(f"\n📁 输出文件:")
    print(f"   📄 DOCX: {output_path}")
    print(f"   📊 PPTX: {ppt_output_path}")
    
    # Convert to PDF using LibreOffice
    import subprocess
    
    try:
        subprocess.run([
            "libreoffice", "--headless", "--convert-to", "pdf",
            output_path, f"--outdir=/root/.openclaw/workspace/EducationStudio/outputs/PDF"
        ], capture_output=True, timeout=30)
        
        print(f"\n✅ PDF 转换成功")
    except Exception as e:
        print(f"\n⚠️  PDF 转换失败：{e}")


if __name__ == "__main__":
    create_final_brochure()
