#!/usr/bin/env python3
"""
创建带图标的最终版宣传册（混合配图策略）
- 实景图：从学校官网下载
- ComfyUI 生成图：数据可视化、图标等
"""

from docx import Document
from pptx import Presentation
from pathlib import Path


def create_docx_with_images():
    """创建带图像的 DOCX 宣传册"""
    
    print("=" * 60)
    print("创建带图像版 DOCX")
    print("=" * 60)
    
    doc = Document()
    
    # Slide 1: Cover with official school photo
    title = "中华书院国际教育合作项目说明书"
    subtitle = "华侨生 985/211 升学项目（校长版）"
    
    doc.add_heading(title, level=1)
    
    p = doc.add_paragraph()
    run = p.add_run(subtitle)
    run.font.size = 14
    
    # Add official school image (placeholder if not downloaded yet)
    ppcha_logo = "/root/.openclaw/workspace/EducationStudio/assets/official_images/ppcha_school.jpg"
    
    # Create cover with both logo and ComfyUI-generated background
    from docx.shared import Inches, Pt
    
    try:
        # 添加 PPCHA 学校图片（如果存在）
        if Path(ppcha_logo).exists():
            logo = doc.add_picture(
                ppcha_logo, 
                width=Inches(4)
            )
            print(f"\n✅ 已嵌入：{ppcha_logo}")
            
        else:
            # 使用 ComfyUI 生成的占位图或生成新图像
            print(f"\n⚠️  PPCHA logo 未找到，使用占位符")
            
    except Exception as e:
        print(f"⚠️  添加图像失败：{e}")
    
    # Add content sections with ComfyUI-generated images
    
    # Section: Why Partner? (with infographic chart)
    doc.add_heading("为什么选择中华书院？", level=2)
    
    p = doc.add_paragraph()
    run = p.add_run("""
📊 **三大核心优势：**

1. ✅ 升学保障 - 双向承诺书制度，未录取即退费
2. ✅ 多元路径 - 985/211 + 海外名校双通道  
3. ✅ 文化传承 - 中华文化与国际视野并重""")
    
    run.font.size = Pt(10)
    
    # Add ComfyUI-generated infographic (placeholder)
    try:
        from docx.shared import Inches
        
        # Use ComfyUI image if available, or create placeholder
        comfyui_infographic = "/root/.openclaw/workspace/EducationStudio/assets/comfyui_generated/admission_rate_chart.png"
        
        if Path(comfyui_infographic).exists():
            chart = doc.add_picture(
                comfyui_infographic,
                width=Inches(6)
            )
            print(f"\n✅ 已嵌入：{comfyui_infographic}")
        else:
            # Create text-based chart instead
            p2 = doc.add_paragraph()
            run2 = p2.add_run("""
📈 **录取率数据示例：**

- 优秀生源 → 985/211 (录取率：85%+)  
- 普通生源 → 211 重点大学 (录取率：70%+)  
- 保底方案 → 暨南大学本科直录 (100%)""")
            run2.font.size = Pt(9)
            
    except Exception as e:
        print(f"⚠️  图表添加失败：{e}")
    
    # Save DOCX
    output_path = "/root/.openclaw/workspace/EducationStudio/outputs/中华书院_华侨生保录宣传册_校长版_final.docx"
    doc.save(output_path)
    
    print(f"\n✅ DOCX 已保存到：{output_path}")


def create_ppt_with_images():
    """创建带图像的 PPT 演示文稿"""
    
    prs = Presentation()
    
    # Slide 1: Cover with PPCHA + ComfyUI images
    slide_layout = prs.slide_layouts[0]  # Title only
    slide = prs.slides.add_slide(slide_layout)
    
    title = "中华书院国际教育合作项目说明书"
    subtitle = "华侨生 985/211 升学项目（校长版）"
    
    slide.shapes.title.text = f"{title}\n{subtitle}"
    
    # Add PPCHA logo (if available)
    ppcha_logo = "/root/.openclaw/workspace/EducationStudio/assets/official_images/ppcha_school.jpg"
    
    try:
        from pptx.util import Inches
        
        if Path(ppcha_logo).exists():
            shape = slide.shapes.add_picture(
                ppcha_logo,
                left=Inches(1.0),
                top=Inches(2.5),
                width=Inches(3.0)
            )
            print(f"✅ PPT: 已嵌入 {ppcha_logo}")
        else:
            # Use ComfyUI-generated image instead
            comfyui_cover = "/root/.openclaw/workspace/EducationStudio/assets/comfyui_generated/school_entrance.png"
            
            if Path(comfyui_cover).exists():
                shape = slide.shapes.add_picture(
                    comfyui_cover,
                    left=Inches(0.5),
                    top=Inches(2.0),
                    width=Inches(6.0)
                )
                print(f"✅ PPT: 已嵌入 ComfyUI 封面图 {comfyui_cover}")
            
            else:
                # Create text placeholder
                shape = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(4.0), Inches(3.0))
                tf = shape.text_frame
                p = tf.paragraphs[0]
                p.text = "📷 [封面图占位符]"
                
    except Exception as e:
        print(f"⚠️  PPT 图像添加失败：{e}")
    
    # Slide 2: Why Partner? (with infographic)
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    title_shape.text = "为什么选择中华书院？"
    
    tf = content_shape.text_frame
    
    for point in [
        "- 📉 生源竞争持续加剧",
        "- 🎯 家长需求升级转变",
        "- 💡 发展路径亟待拓宽",
        ""
    ]:
        if point:
            p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
            p.text = point
    
    # Add ComfyUI-generated infographic (if available)
    try:
        comfyui_pain_points = "/root/.openclaw/workspace/EducationStudio/assets/comfyui_generated/industry_chart.png"
        
        if Path(comfyui_pain_points).exists():
            shape = slide.shapes.add_picture(
                comfyui_pain_points,
                left=Inches(7.0),
                top=Inches(1.0),
                width=Inches(5.0)
            )
            
    except Exception as e:
        print(f"⚠️  PPT 图表添加失败：{e}")
    
    # Save PPTX
    output_path = "/root/.openclaw/workspace/EducationStudio/outputs/中华书院_华侨生保录宣传册_校长版_final.pptx"
    prs.save(output_path)
    
    print(f"\n✅ PPTX 已保存到：{output_path}")


def main():
    """主函数"""
    
    # Step 1: Create DOCX with images
    create_docx_with_images()
    
    # Step 2: Create PPT with images  
    create_ppt_with_images()
    
    print("\n" + "=" * 60)
    print("✅ 最终版宣传册创建完成！")
    print("=" * 60)


if __name__ == "__main__":
    main()
