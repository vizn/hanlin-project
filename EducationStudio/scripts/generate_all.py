#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
中华书院宣传册一键生成脚本
支持：DOCX + PDF + PPT 三种格式输出
"""

import subprocess
import sys
from pathlib import Path


def generate_docx():
    """生成 DOCX 格式"""
    from docx import Document
    
    # 读取 Markdown 源文件
    markdown_file = "/root/.openclaw/workspace/中华书院_华侨生保录宣传册_校长版.md"
    
    try:
        with open(markdown_file, 'r', encoding='utf-8') as f:
            md_content = f.read()
        
        doc = Document()
        
        # 添加标题
        doc.add_heading('中华书院国际教育合作项目说明书', level=1)
        doc.add_heading('华侨生 985/211 升学项目（校长版）', level=2)
        
        # 分割内容并添加到文档
        sections = md_content.split('### ')
        
        for i, section in enumerate(sections[1:], 1):  # 跳过 P1 封面部分
            if 'P' in section:
                lines = section.strip().split('\n')
                
                # 提取页码和标题
                header_line = ""
                content_lines = []
                
                for line in lines:
                    if line.startswith('### P'):
                        page_num, title = line.split(' ', 1)
                        doc.add_heading(f"{page_num} {title}", level=2)
                        header_line = f"P{page_num}: {title}"
                    else:
                        content_lines.append(line.strip())
                
                # 添加内容
                for line in content_lines:
                    if line.strip():
                        doc.add_paragraph(line.strip(), style='List Bullet')
        
        # 保存 DOCX
        output_path = "/root/.openclaw/workspace/EducationStudio/outputs/中华书院_华侨生保录宣传册_校长版.docx"
        doc.save(output_path)
        
        print(f"✅ DOCX 生成成功：{output_path}")
        return True
        
    except Exception as e:
        print(f"❌ DOCX 生成失败：{e}")
        return False


def generate_pdf():
    """生成 PDF 格式（使用 pandoc + wkhtmltopdf）"""
    
    try:
        # 先尝试用 pdflatex（如果有 LaTeX）
        subprocess.run(['pdflatex', '-interaction=nonstopmode', 
                       '/root/.openclaw/workspace/中华书院_华侨生保录宣传册_校长版.md',
                       '-output-directory=/root/.openclaw/workspace/EducationStudio/outputs'],
                      capture_output=True, text=True)
        
        output_path = "/root/.openclaw/workspace/EducationStudio/outputs/中华书院_华侨生保录宣传册_校长版.pdf"
        
        if Path(output_path).exists():
            print(f"✅ PDF 生成成功：{output_path}")
            return True
        else:
            # 尝试用 wkhtmltopdf
            subprocess.run(['wkhtmltohtml', 
                          '/root/.openclaw/workspace/中华书院_华侨生保录宣传册_校长版.md',
                          '-o', output_path.replace('.pdf', '.html')],
                         capture_output=True, text=True)
            
            # 如果 wkhtmltohtml 也不可用，用简单方法
            print("⚠️  PDF 生成依赖未安装，正在使用备用方案...")
            
    except FileNotFoundError:
        pass
    
    return False


def generate_ppt():
    """生成 PPTX 格式"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    
    try:
        prs = Presentation()
        
        # 读取 Markdown 内容
        with open('/root/.openclaw/workspace/中华书院_华侨生保录宣传册_校长版.md', 'r') as f:
            content = f.read()
        
        # P1：封面
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "中华书院国际教育合作项目说明书"
        subtitle.text = "华侨生 985/211 升学项目（校长版）\n\n打造学校国际化特色 · 提升招生竞争力 · 构建多元升学通道"
        
        # P2-P7：内容页
        for page in content.split('### ')[1:]:  # 跳过第一部分
            if 'P' not in page or not page.strip():
                continue
                
            lines = page.strip().split('\n')
            
            # 提取标题和内容
            title_line = ""
            bullets = []
            
            for line in lines:
                if line.startswith('## ') or (line and not line.startswith('-')):
                    if '### P' in page:
                        # 查找页码和标题
                        parts = [p.strip() for p in page.split('\n')]
                        for part in parts:
                            if part.startswith('### P'):
                                title_line = part.replace('### P', '').strip()
                            elif part and not part.startswith('-') and '##' not in part:
                                title_line = part.replace('———', '').replace('*', '')
                            
            # 简化处理：直接创建内容页
            slide_layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            
            if lines:
                title_shape = slide.shapes.title
                content_shape = slide.placeholders[1]
                
                # 提取当前页标题
                current_title = ""
                for line in lines:
                    if '### P' not in line and '##' not in line and line.strip():
                        current_title = line.strip().replace('*', '').strip()
                
                title_shape.text = current_title
                
                tf = content_shape.text_frame
                tf.text = ""
                
        # 保存 PPTX
        output_path = "/root/.openclaw/workspace/EducationStudio/outputs/中华书院_华侨生保录宣传册_校长版.pptx"
        prs.save(output_path)
        
        print(f"✅ PPTX 生成成功：{output_path}")
        return True
        
    except Exception as e:
        print(f"❌ PPTX 生成失败：{e}")
        return False


def main():
    """主函数"""
    print("=" * 60)
    print("中华书院宣传册一键生成工具")
    print("=" * 60)
    
    # 检查依赖
    required_deps = ['python-docx', 'pptx']
    missing_deps = []
    
    for dep in required_deps:
        try:
            __import__(dep.replace('-', '_'))
        except ImportError:
            missing_deps.append(dep)
    
    if missing_deps:
        print(f"⚠️  缺少依赖：{', '.join(missing_deps)}")
        print("正在安装...")
        subprocess.run(['pip3', 'install'] + missing_deps, capture_output=True)
        
        # 重新导入
        for dep in required_deps:
            try:
                __import__(dep.replace('-', '_'))
            except ImportError:
                pass
    
    # 生成各格式
    print("\n📄 正在生成 DOCX...")
    generate_docx()
    
    print("\n📑 正在生成 PPTX...")
    generate_ppt()
    
    # PDF（如果有 pandoc/wkhtmltopdf）
    try:
        import subprocess
        result = subprocess.run(
            ['pandoc', '-t', 'pdf', 
             '/root/.openclaw/workspace/中华书院_华侨生保录宣传册_校长版.md',
             '-o', '/root/.openclaw/workspace/EducationStudio/outputs/中华书院_华侨生保录宣传册_校长版.pdf'],
            capture_output=True, text=True
        )
        
        if result.returncode == 0:
            print("\n📕 PDF 生成成功：pandoc")
        else:
            # 备用方案：使用 LibreOffice
            subprocess.run(
                ['libreoffice', '--headless', '--convert-to', 'pdf',
                 '/root/.openclaw/workspace/中华书院_华侨生保录宣传册_校长版.docx'],
                capture_output=True, text=True
            )
            
    except FileNotFoundError:
        print("\n⚠️  PDF 生成依赖不可用，跳过 PDF 生成")
    
    # 列出输出文件
    output_dir = "/root/.openclaw/workspace/EducationStudio/outputs"
    print("\n" + "=" * 60)
    print("✅ 生成完成！输出文件列表：")
    print("=" * 60)
    
    for file in Path(output_dir).glob("*"):
        size = file.stat().st_size / 1024
        print(f"📁 {file.name} ({size:.1f} KB)")
    
    print("\n提示：如需 PDF，请确保安装了 pandoc 或 libreoffice")


if __name__ == "__main__":
    main()
