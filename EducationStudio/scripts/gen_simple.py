#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""简化版一键生成脚本"""

from docx import Document
from pptx import Presentation
import os


def generate_docx():
    """生成 DOCX"""
    md_file = "/root/.openclaw/workspace/中华书院_华侨生保录宣传册_校长版.md"
    
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    doc = Document()
    doc.add_heading('中华书院国际教育合作项目说明书', 0)
    doc.add_heading('华侨生 985/211 升学项目（校长版）', level=1)
    
    # 添加内容
    for line in content.split('\n'):
        if line.strip() and not line.startswith('#') and '### P' not in line:
            doc.add_paragraph(line.strip(), style='List Bullet')
    
    output = "/root/.openclaw/workspace/EducationStudio/outputs/中华书院_华侨生保录宣传册_校长版.docx"
    doc.save(output)
    print(f"✅ DOCX 生成：{output}")


def generate_ppt():
    """生成 PPT"""
    md_file = "/root/.openclaw/workspace/中华书院_华侨生保录宣传册_校长版.md"
    
    prs = Presentation()
    
    # 封面
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "中华书院国际教育合作项目说明书"
    subtitle = slide.placeholders[1].text_frame.paragraphs[0]
    subtitle.text = "华侨生 985/211 升学项目（校长版）\n打造学校国际化特色 · 提升招生竞争力 · 构建多元升学通道"
    
    # P2-P8：内容页
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    for page in content.split('### ')[1:]:
        if not page.strip():
            continue
        
        lines = [l.strip() for l in page.split('\n') if l.strip()]
        
        # 提取第一行作为标题（非列表项）
        title = ""
        bullets = []
        
        for line in lines:
            if '### P' not in line and line and not line.startswith('-') and '*' not in line:
                if '##' not in line or 'P' in page:
                    title = line.strip().replace('*', '').strip()
            elif line.startswith('-') or line.startswith('•'):
                bullets.append(line.strip())
        
        # 如果提取到内容，创建新页
        if lines and ('### P' in page or '##' not in lines[0]):
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            
            title_shape = slide.shapes.title
            content_shape = slide.placeholders[1]
            
            # 简化：只取第一行非列表内容作为标题
            current_title = ""
            for line in lines:
                if '### P' not in line and line.strip() and '*' not in line and '-' not in line[0]:
                    current_title = line.strip().replace('——', '').strip()
            
            title_shape.text = current_title
            
            tf = content_shape.text_frame
            for bullet in bullets[:15]:  # 限制每页行数
                p = tf.add_paragraph()
                p.text = bullet.replace('*', '')
                p.level = 0
    
    output = "/root/.openclaw/workspace/EducationStudio/outputs/中华书院_华侨生保录宣传册_校长版.pptx"
    prs.save(output)
    print(f"✅ PPTX 生成：{output}")


def main():
    print("=" * 60)
    print("中华书院宣传册一键生成")
    print("=" * 60)
    
    generate_docx()
    generate_ppt()
    
    # 列出输出文件
    output_dir = "/root/.openclaw/workspace/EducationStudio/outputs"
    print("\n📁 输出文件：")
    for f in os.listdir(output_dir):
        path = os.path.join(output_dir, f)
        size = os.path.getsize(path) / 1024
        print(f"   📄 {f} ({size:.1f} KB)")


if __name__ == "__main__":
    main()
